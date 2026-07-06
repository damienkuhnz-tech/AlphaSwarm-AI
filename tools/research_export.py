"""
╔══════════════════════════════════════════════════════════════════════════════╗
║  TOOLS / RESEARCH EXPORT — GÉNÉRATION DE LIVRABLES MULTI-FORMATS            ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  Génère Excel (.xlsx), Word (.docx), PowerPoint (.pptx) à partir des        ║
║  analyses Research et du brief de cadrage. Stocké dans                      ║
║  outputs/research_exports/.                                                  ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  Toutes les fonctions ont la même signature :                               ║
║    fn(research: list[dict], brief: dict, out_dir: str) -> str (path)       ║
║  research = liste de ResearchOutput sérialisés (dict).                       ║
║  brief    = brief structuré du briefing PM (peut être vide).                 ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

# ── Imports ───────────────────────────────────────────────────────────────────
import os
import re
from datetime import date
from typing import Any, Dict, List, Optional, Tuple


# ── Helpers communs ───────────────────────────────────────────────────────────

def _safe_get(d: Any, *keys, default=None):
    """Lecture sécurisée dans un dict imbriqué (et sur attributs Pydantic)."""
    cur = d
    for k in keys:
        if cur is None:
            return default
        if isinstance(cur, dict):
            cur = cur.get(k)
        else:
            cur = getattr(cur, k, None)
    return cur if cur is not None else default


def _brief_label(brief: Dict[str, Any]) -> str:
    """Label court pour titres de documents : 'Long Quality Tech Mega US' etc."""
    parts = []
    for key in ("long_short", "style", "univers", "sous_secteur", "taille_capi", "geo"):
        v = brief.get(key)
        if v and isinstance(v, str):
            parts.append(v.replace("_", " ").title())
    return " ".join(parts) or "Equity Research"


def _today_stamp() -> str:
    return date.today().strftime("%Y%m%d")


# ╔═════════════════════════════════════════════════════════════════════════════╗
# ║  PARSER MARKDOWN BASIQUE                                                     ║
# ║  Convertit un markdown en blocs structurés : (type, level, content)         ║
# ║  Types : "h1", "h2", "h3", "para", "bullet", "table"                       ║
# ╚═════════════════════════════════════════════════════════════════════════════╝

def _parse_markdown(md: str) -> List[Tuple[str, int, Any]]:
    """
    Parse markdown ligne par ligne en blocs (type, level, content).
    Suffisamment robuste pour les sorties LLM (pas un parser exhaustif).
    """
    lines = md.split("\n")
    blocks: List[Tuple[str, int, Any]] = []
    i = 0
    n = len(lines)
    para_buf: List[str] = []
    bullet_buf: List[str] = []

    def _flush_para():
        if para_buf:
            text = " ".join(para_buf).strip()
            if text:
                blocks.append(("para", 0, text))
            para_buf.clear()

    def _flush_bullets():
        if bullet_buf:
            blocks.append(("bullet", 0, list(bullet_buf)))
            bullet_buf.clear()

    while i < n:
        line = lines[i].rstrip()
        stripped = line.strip()

        # ── Headings ──────────────────────────────────────────────────────────
        m = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if m:
            _flush_para(); _flush_bullets()
            level = len(m.group(1))
            blocks.append(("h" + str(min(level, 3)), level, m.group(2).strip()))
            i += 1
            continue

        # ── Tables markdown (| col1 | col2 | ... |) ───────────────────────────
        if stripped.startswith("|") and "|" in stripped[1:]:
            table_rows = []
            while i < n and lines[i].strip().startswith("|"):
                row = lines[i].strip()
                # Skip separator row (---)
                if not re.match(r"^\|[\s\-:|]+\|$", row):
                    cells = [c.strip() for c in row.strip("|").split("|")]
                    table_rows.append(cells)
                i += 1
            if table_rows:
                _flush_para(); _flush_bullets()
                blocks.append(("table", 0, table_rows))
            continue

        # ── Bullets ───────────────────────────────────────────────────────────
        if re.match(r"^[\-\*\+]\s+", stripped) or re.match(r"^\d+[\.\)]\s+", stripped):
            _flush_para()
            content = re.sub(r"^[\-\*\+]\s+|^\d+[\.\)]\s+", "", stripped)
            bullet_buf.append(content)
            i += 1
            continue

        # ── Ligne vide → flush ───────────────────────────────────────────────
        if not stripped:
            _flush_para(); _flush_bullets()
            i += 1
            continue

        # ── Paragraphe ───────────────────────────────────────────────────────
        _flush_bullets()
        para_buf.append(stripped)
        i += 1

    _flush_para(); _flush_bullets()
    return blocks


def _strip_markdown_inline(text: str) -> str:
    """Retire les marqueurs markdown inline (gras, italique, code) pour rendu propre."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"__(.+?)__",     r"\1", text)
    text = re.sub(r"\*(.+?)\*",     r"\1", text)
    text = re.sub(r"_(.+?)_",       r"\1", text)
    text = re.sub(r"`(.+?)`",       r"\1", text)
    return text


# ╔═════════════════════════════════════════════════════════════════════════════╗
# ║  EXPORT EXCEL                                                                ║
# ║  Workbook avec feuilles : Idées | Détails | Comparables | DCF | Brief       ║
# ╚═════════════════════════════════════════════════════════════════════════════╝

def export_research_to_excel(
    research: List[Dict[str, Any]],
    brief:    Dict[str, Any],
    out_dir:  str,
) -> str:
    """Génère un .xlsx multi-feuilles avec les analyses."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    wb.remove(wb.active)  # Supprime la feuille par défaut

    # Style header commun
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill("solid", fgColor="0E4A2C")
    thin_border = Border(left=Side(style="thin", color="DDDDDD"),
                         right=Side(style="thin", color="DDDDDD"),
                         top=Side(style="thin", color="DDDDDD"),
                         bottom=Side(style="thin", color="DDDDDD"))

    def _format_header_row(ws, row_num=1):
        for cell in ws[row_num]:
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

    # ── Feuille 1 : Idées (synthèse) ──────────────────────────────────────────
    ws = wb.create_sheet("Idées")
    ws.append(["Ticker", "Nom", "Secteur", "Géo", "Rating", "Score conviction",
               "Prix actuel", "Prix cible bas", "Prix cible moyen", "Prix cible haut",
               "Upside %", "Poids suggéré", "Recommandation", "Thèse résumée"])
    _format_header_row(ws)

    for r in research:
        valuation = _safe_get(r, "valorisation") or {}
        prix_act = _safe_get(r, "prix_actuel") or _safe_get(valuation, "prix_actuel")
        cible    = _safe_get(valuation, "prix_cible_central") or _safe_get(valuation, "prix_cible_moyen")
        upside = None
        if prix_act and cible:
            try:
                upside = round((float(cible) - float(prix_act)) / float(prix_act) * 100, 1)
            except (TypeError, ValueError):
                pass
        these = _safe_get(r, "these_investissement") or ""
        ws.append([
            _safe_get(r, "ticker"),
            _safe_get(r, "nom"),
            _safe_get(r, "secteur"),
            _safe_get(r, "geographie"),
            _safe_get(r, "recommandation") or _safe_get(r, "rating"),
            _safe_get(r, "score_conviction"),
            prix_act,
            _safe_get(valuation, "prix_cible_bas"),
            cible,
            _safe_get(valuation, "prix_cible_haut"),
            upside,
            _safe_get(r, "poids_suggere_initial"),
            _safe_get(r, "recommandation"),
            (these[:300] + "...") if len(these) > 300 else these,
        ])
    # Largeurs auto
    for col in ws.columns:
        max_len = max((len(str(c.value or "")) for c in col), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(max(max_len + 2, 10), 60)

    # ── Feuille 2 : Détails par titre (1 bloc par ticker) ─────────────────────
    ws = wb.create_sheet("Détails")
    row = 1
    for r in research:
        ticker = _safe_get(r, "ticker") or "?"
        ws.cell(row=row, column=1, value=f"{ticker} — {_safe_get(r, 'nom') or ''}").font = Font(bold=True, size=14, color="0E4A2C")
        row += 1
        ws.cell(row=row, column=1, value=f"Recommandation : {_safe_get(r, 'recommandation') or 'N/D'}    "
                f"| Score : {_safe_get(r, 'score_conviction') or 'N/D'}/100").font = Font(italic=True)
        row += 2

        sections = [
            ("Thèse d'investissement", _safe_get(r, "these_investissement")),
            ("Catalyseurs",            ", ".join(_safe_get(r, "catalyseurs", default=[]) or [])),
            ("Risques principaux",     ", ".join(_safe_get(r, "risques_principaux", default=[]) or [])),
            ("Avantages compétitifs",  ", ".join(_safe_get(r, "avantages_competitifs", default=[]) or [])),
            ("Conclusion",             _safe_get(r, "conclusion") or _safe_get(r, "executive_summary")),
        ]
        for title, content in sections:
            if not content:
                continue
            ws.cell(row=row, column=1, value=title).font = Font(bold=True, color="0E4A2C")
            row += 1
            ws.cell(row=row, column=1, value=str(content)[:1000]).alignment = Alignment(wrap_text=True, vertical="top")
            ws.row_dimensions[row].height = 60
            row += 2
        row += 2  # Saut entre tickers

    ws.column_dimensions["A"].width = 120

    # ── Feuille 3 : Comparables (ratios fondamentaux) ─────────────────────────
    ws = wb.create_sheet("Comparables")
    ws.append(["Ticker", "Nom", "P/E TTM", "P/E NTM", "EV/EBITDA",
               "P/B", "P/S", "ROE", "Marge op.", "Marge nette", "Dette/Capitaux",
               "Croissance revenus YoY", "Beta"])
    _format_header_row(ws)
    for r in research:
        f = _safe_get(r, "fundamentaux") or _safe_get(r, "valorisation") or {}
        ws.append([
            _safe_get(r, "ticker"),
            _safe_get(r, "nom"),
            _safe_get(f, "PER_ttm")  or _safe_get(f, "pe_ttm"),
            _safe_get(f, "PER_forward") or _safe_get(f, "pe_forward") or _safe_get(f, "PER_estime_NTM"),
            _safe_get(f, "EV_EBITDA") or _safe_get(f, "ev_ebitda"),
            _safe_get(f, "PB_ratio") or _safe_get(f, "pb_ratio"),
            _safe_get(f, "PS_ratio") or _safe_get(f, "ps_ratio"),
            _safe_get(f, "ROE") or _safe_get(f, "roe"),
            _safe_get(f, "marge_operationnelle"),
            _safe_get(f, "marge_nette"),
            _safe_get(f, "dette_sur_capitaux"),
            _safe_get(f, "croissance_revenus_yoy"),
            _safe_get(f, "beta"),
        ])
    for col in ws.columns:
        max_len = max((len(str(c.value or "")) for c in col), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(max(max_len + 2, 12), 30)

    # ── Feuille 4 : Brief (paramètres de la recherche) ────────────────────────
    ws = wb.create_sheet("Brief")
    ws.append(["Paramètre", "Valeur"])
    _format_header_row(ws)
    for k, v in (brief or {}).items():
        if isinstance(v, (list, dict)):
            v = str(v)
        ws.append([k, v])
    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 60

    # ── Sauvegarde ────────────────────────────────────────────────────────────
    label = _brief_label(brief).replace(" ", "_")[:40]
    fname = f"Research_{label}_{_today_stamp()}.xlsx"
    path = os.path.join(out_dir, fname)
    wb.save(path)
    return path


# ╔═════════════════════════════════════════════════════════════════════════════╗
# ║  EXPORT WORD                                                                 ║
# ║  Document narratif : couverture + synthèse + 1 section par ticker           ║
# ╚═════════════════════════════════════════════════════════════════════════════╝

def export_research_to_word(
    research: List[Dict[str, Any]],
    brief:    Dict[str, Any],
    out_dir:  str,
) -> str:
    """Génère un .docx narratif avec l'analyse complète."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # ── Style global ──────────────────────────────────────────────────────────
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # ── Page de garde ─────────────────────────────────────────────────────────
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("EQUITY RESEARCH REPORT")
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0E, 0x4A, 0x2C)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub.add_run(_brief_label(brief))
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(date.today().strftime("%d %B %Y"))
    run.font.size = Pt(11)
    run.font.italic = True

    doc.add_page_break()

    # ── Résumé du brief ───────────────────────────────────────────────────────
    doc.add_heading("Cadrage de la recherche", level=1)
    if brief:
        for k, v in brief.items():
            if v in (None, "", [], {}):
                continue
            p = doc.add_paragraph()
            p.add_run(f"{k.replace('_', ' ').title()} : ").bold = True
            p.add_run(str(v))
    else:
        doc.add_paragraph("Aucun brief structuré — recherche standard.")

    doc.add_paragraph()

    # ── Synthèse exécutive ────────────────────────────────────────────────────
    doc.add_heading("Synthèse exécutive", level=1)
    doc.add_paragraph(
        f"Cette recherche couvre {len(research)} idées d'investissement. "
        f"Les sections suivantes détaillent la thèse, les catalyseurs, les risques "
        f"et la valorisation pour chaque titre."
    )

    # Tableau récap
    table = doc.add_table(rows=1, cols=5)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Ticker"
    hdr[1].text = "Nom"
    hdr[2].text = "Recommandation"
    hdr[3].text = "Score"
    hdr[4].text = "Upside %"
    for r in research:
        row = table.add_row().cells
        row[0].text = str(_safe_get(r, "ticker") or "")
        row[1].text = str(_safe_get(r, "nom") or "")[:30]
        row[2].text = str(_safe_get(r, "recommandation") or "N/D")
        row[3].text = str(_safe_get(r, "score_conviction") or "N/D")
        prix = _safe_get(r, "prix_actuel")
        cible = _safe_get(r, "valorisation", "prix_cible_central") or _safe_get(r, "valorisation", "prix_cible_moyen")
        if prix and cible:
            try:
                row[4].text = f"{round((float(cible) - float(prix)) / float(prix) * 100, 1)}%"
            except Exception:
                row[4].text = "N/D"
        else:
            row[4].text = "N/D"

    doc.add_page_break()

    # ── Section par ticker ────────────────────────────────────────────────────
    for r in research:
        ticker = _safe_get(r, "ticker") or "?"
        nom    = _safe_get(r, "nom") or ""
        doc.add_heading(f"{ticker} — {nom}", level=1)

        meta = doc.add_paragraph()
        meta.add_run(f"Recommandation : {_safe_get(r, 'recommandation') or 'N/D'}    ").bold = True
        meta.add_run(f"Score conviction : {_safe_get(r, 'score_conviction') or 'N/D'}/100    ")
        meta.add_run(f"Secteur : {_safe_get(r, 'secteur') or 'N/D'}    ")
        meta.add_run(f"Géo : {_safe_get(r, 'geographie') or 'N/D'}")

        sections = [
            ("Thèse d'investissement",     _safe_get(r, "these_investissement")),
            ("Catalyseurs",                _safe_get(r, "catalyseurs", default=[])),
            ("Avantages compétitifs",      _safe_get(r, "avantages_competitifs", default=[])),
            ("Risques principaux",         _safe_get(r, "risques_principaux", default=[])),
            ("Conclusion",                 _safe_get(r, "conclusion") or _safe_get(r, "executive_summary")),
        ]
        for title, content in sections:
            if not content:
                continue
            doc.add_heading(title, level=2)
            if isinstance(content, list):
                for item in content:
                    doc.add_paragraph(str(item), style="List Bullet")
            else:
                doc.add_paragraph(str(content))

        doc.add_page_break()

    # ── Sauvegarde ────────────────────────────────────────────────────────────
    label = _brief_label(brief).replace(" ", "_")[:40]
    fname = f"Research_{label}_{_today_stamp()}.docx"
    path = os.path.join(out_dir, fname)
    doc.save(path)
    return path


# ╔═════════════════════════════════════════════════════════════════════════════╗
# ║  EXPORT POWERPOINT                                                           ║
# ║  Pitch deck : Cover | Brief | Synthèse | 1 slide / idée | Conclusion        ║
# ╚═════════════════════════════════════════════════════════════════════════════╝

def export_research_to_pptx(
    research: List[Dict[str, Any]],
    brief:    Dict[str, Any],
    out_dir:  str,
) -> str:
    """Génère un .pptx pitch deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_DARK   = RGBColor(0x0A, 0x14, 0x18)
    ACCENT    = RGBColor(0x22, 0xC5, 0x5E)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    GREY      = RGBColor(0xAA, 0xAA, 0xAA)

    def _add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        from pptx.util import Inches
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tf

    def _set_bg(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    blank = prs.slide_layouts[6]

    # ── Slide 1 : Cover ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text_box(s, 0.5, 2.5, 12.3, 1.5, "EQUITY RESEARCH",
                  size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 4.0, 12.3, 1.0, _brief_label(brief),
                  size=24, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 6.5, 12.3, 0.5, date.today().strftime("%d %B %Y"),
                  size=14, color=GREY, align=PP_ALIGN.CENTER)

    # ── Slide 2 : Cadrage ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text_box(s, 0.5, 0.4, 12.3, 0.7, "Cadrage de la recherche",
                  size=28, bold=True, color=ACCENT)
    y = 1.4
    for k, v in (brief or {}).items():
        if v in (None, "", [], {}):
            continue
        _add_text_box(s, 0.7, y, 12.0, 0.4,
                      f"{k.replace('_', ' ').title()} : {v}", size=14, color=WHITE)
        y += 0.45
        if y > 6.5:
            break

    # ── Slide 3 : Synthèse (tableau récap) ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text_box(s, 0.5, 0.4, 12.3, 0.7, f"Synthèse — {len(research)} idées",
                  size=28, bold=True, color=ACCENT)
    rows = len(research) + 1
    table_shape = s.shapes.add_table(rows, 5, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4 * rows))
    tbl = table_shape.table
    headers = ["Ticker", "Nom", "Reco", "Score", "Upside %"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(13)
    for r_idx, r in enumerate(research, start=1):
        prix = _safe_get(r, "prix_actuel")
        cible = _safe_get(r, "valorisation", "prix_cible_central") or _safe_get(r, "valorisation", "prix_cible_moyen")
        upside = "N/D"
        if prix and cible:
            try:
                upside = f"{round((float(cible) - float(prix)) / float(prix) * 100, 1)}%"
            except Exception:
                pass
        vals = [
            str(_safe_get(r, "ticker") or ""),
            str(_safe_get(r, "nom") or "")[:24],
            str(_safe_get(r, "recommandation") or "N/D"),
            str(_safe_get(r, "score_conviction") or "N/D"),
            upside,
        ]
        for i, v in enumerate(vals):
            cell = tbl.cell(r_idx, i)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)

    # ── Slide par idée ────────────────────────────────────────────────────────
    for r in research:
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        ticker = _safe_get(r, "ticker") or "?"
        nom    = _safe_get(r, "nom") or ""
        _add_text_box(s, 0.5, 0.4, 12.3, 0.7, f"{ticker} — {nom}",
                      size=28, bold=True, color=ACCENT)

        reco = _safe_get(r, "recommandation") or "N/D"
        score = _safe_get(r, "score_conviction") or "N/D"
        _add_text_box(s, 0.5, 1.2, 12.3, 0.4,
                      f"Recommandation : {reco}    Score : {score}/100    "
                      f"Secteur : {_safe_get(r, 'secteur') or 'N/D'}    "
                      f"Géo : {_safe_get(r, 'geographie') or 'N/D'}",
                      size=12, color=GREY)

        # Thèse
        these = _safe_get(r, "these_investissement") or ""
        _add_text_box(s, 0.5, 1.9, 12.3, 0.4, "Thèse", size=16, bold=True, color=ACCENT)
        _add_text_box(s, 0.5, 2.4, 12.3, 1.6, these[:600], size=12, color=WHITE)

        # Catalyseurs (3 max)
        cats = _safe_get(r, "catalyseurs", default=[]) or []
        _add_text_box(s, 0.5, 4.1, 6.0, 0.4, "Catalyseurs", size=16, bold=True, color=ACCENT)
        cats_text = "\n".join(f"• {c}" for c in cats[:4])
        _add_text_box(s, 0.5, 4.6, 6.0, 2.3, cats_text or "—", size=12, color=WHITE)

        # Risques
        risks = _safe_get(r, "risques_principaux", default=[]) or []
        _add_text_box(s, 6.7, 4.1, 6.1, 0.4, "Risques", size=16, bold=True, color=ACCENT)
        risks_text = "\n".join(f"• {r_}" for r_ in risks[:4])
        _add_text_box(s, 6.7, 4.6, 6.1, 2.3, risks_text or "—", size=12, color=WHITE)

    # ── Slide finale : Recommandations ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text_box(s, 0.5, 0.4, 12.3, 0.7, "Recommandations",
                  size=28, bold=True, color=ACCENT)
    buys = [r for r in research if (_safe_get(r, "recommandation") or "").upper() in ("BUY", "STRONG BUY")]
    holds = [r for r in research if (_safe_get(r, "recommandation") or "").upper() in ("HOLD",)]
    sells = [r for r in research if (_safe_get(r, "recommandation") or "").upper() in ("SELL", "STRONG SELL")]
    summary = (
        f"Idées BUY/STRONG BUY : {len(buys)}\n"
        f"Idées HOLD : {len(holds)}\n"
        f"Idées SELL/STRONG SELL : {len(sells)}\n\n"
        f"Top BUY (par score) :"
    )
    _add_text_box(s, 0.5, 1.4, 12.3, 1.5, summary, size=14, color=WHITE)
    top_buys = sorted(buys, key=lambda x: _safe_get(x, "score_conviction", default=0) or 0, reverse=True)[:5]
    bullets = "\n".join(
        f"• {_safe_get(b, 'ticker')} — {_safe_get(b, 'nom')} — Score {_safe_get(b, 'score_conviction')}/100"
        for b in top_buys
    )
    _add_text_box(s, 0.7, 3.5, 12.0, 3.0, bullets or "Aucune idée BUY", size=14, color=ACCENT)

    # ── Sauvegarde ────────────────────────────────────────────────────────────
    label = _brief_label(brief).replace(" ", "_")[:40]
    fname = f"Research_{label}_{_today_stamp()}.pptx"
    path = os.path.join(out_dir, fname)
    prs.save(path)
    return path


# ╔═════════════════════════════════════════════════════════════════════════════╗
# ║  EXPORTS À PARTIR DE MARKDOWN (skills equity-research-like)                 ║
# ║  Convertit le markdown produit par un skill en Excel, Word, PowerPoint.    ║
# ╚═════════════════════════════════════════════════════════════════════════════╝

def _skill_filename(skill_meta: Dict[str, Any], ext: str) -> str:
    """Construit un nom de fichier propre : Earnings_AAPL_Q2_FY26_20260507.xlsx"""
    prefix = skill_meta.get("filename_prefix", "Research")
    answers = skill_meta.get("answers", {}) or {}
    parts = [prefix]
    # Ajoute jusqu'à 3 valeurs des réponses pour disambiguer
    for key in ("ticker", "tickers", "trimestre", "annee_fiscale", "univers", "secteur",
                "style", "geo", "side", "horizon"):
        v = answers.get(key)
        if v and isinstance(v, (str, int)):
            s = str(v).replace(" ", "").replace("/", "-").replace(",", "_")
            if len(s) > 0 and len(s) <= 20 and len(parts) < 4:
                parts.append(s)
    parts.append(_today_stamp())
    name = "_".join(parts)
    name = re.sub(r"[^A-Za-z0-9_\-]", "", name)[:80]
    return f"{name}.{ext}"


def export_skill_to_excel(markdown: str, skill_meta: Dict[str, Any], out_dir: str) -> str:
    """
    Génère un .xlsx à partir du markdown d'un skill.
    Feuilles :
      - Synthèse : sections H1/H2 du markdown
      - Tableaux : tous les tableaux trouvés dans le markdown (1 par feuille si gros)
      - Brief : paramètres ayant servi à générer
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    blocks = _parse_markdown(markdown)

    wb = Workbook()
    wb.remove(wb.active)

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill("solid", fgColor="0E4A2C")
    h1_font = Font(bold=True, size=14, color="0E4A2C")
    h2_font = Font(bold=True, size=12, color="0E4A2C")

    # ── Feuille 1 : Synthèse ──────────────────────────────────────────────────
    ws = wb.create_sheet("Synthèse")
    ws.column_dimensions["A"].width = 110
    row = 1
    title = skill_meta.get("label", "Research") + " — " + date.today().strftime("%d %B %Y")
    ws.cell(row=row, column=1, value=title).font = Font(bold=True, size=18, color="0E4A2C")
    row += 2

    for btype, level, content in blocks:
        if btype == "h1" or (btype == "h2" and level == 1):
            cell = ws.cell(row=row, column=1, value=_strip_markdown_inline(content))
            cell.font = h1_font
            row += 1
        elif btype == "h2":
            cell = ws.cell(row=row, column=1, value=_strip_markdown_inline(content))
            cell.font = h2_font
            row += 1
        elif btype == "h3":
            cell = ws.cell(row=row, column=1, value=_strip_markdown_inline(content))
            cell.font = Font(bold=True, size=11)
            row += 1
        elif btype == "para":
            cell = ws.cell(row=row, column=1, value=_strip_markdown_inline(content))
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            ws.row_dimensions[row].height = max(15, min(60, len(content) // 2))
            row += 1
        elif btype == "bullet":
            for item in content:
                cell = ws.cell(row=row, column=1, value="• " + _strip_markdown_inline(item))
                cell.alignment = Alignment(wrap_text=True)
                row += 1
        elif btype == "table":
            # Saute, on les met dans Tableaux
            row += 1
            continue
        row += 0  # spacer

    # ── Feuille 2 : Tableaux ──────────────────────────────────────────────────
    tables = [c for (t, _, c) in blocks if t == "table"]
    if tables:
        ws_t = wb.create_sheet("Tableaux")
        cur_row = 1
        for idx, table in enumerate(tables, 1):
            # Titre du tableau
            ws_t.cell(row=cur_row, column=1, value=f"Tableau {idx}").font = Font(bold=True, size=12, color="0E4A2C")
            cur_row += 1
            # Header
            if table:
                for j, cell_val in enumerate(table[0], start=1):
                    c = ws_t.cell(row=cur_row, column=j, value=_strip_markdown_inline(cell_val))
                    c.font = header_font
                    c.fill = header_fill
                    c.alignment = Alignment(horizontal="center")
                cur_row += 1
                # Data rows
                for tr in table[1:]:
                    for j, cell_val in enumerate(tr, start=1):
                        ws_t.cell(row=cur_row, column=j, value=_strip_markdown_inline(cell_val))
                    cur_row += 1
            cur_row += 2

        # Largeurs auto
        for col in ws_t.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=12)
            ws_t.column_dimensions[col[0].column_letter].width = min(max(max_len + 2, 12), 50)

    # ── Feuille 3 : Brief (paramètres) ────────────────────────────────────────
    ws_b = wb.create_sheet("Brief")
    ws_b.append(["Paramètre", "Valeur"])
    for cell in ws_b[1]:
        cell.font = header_font
        cell.fill = header_fill
    answers = skill_meta.get("answers", {}) or {}
    for k, v in answers.items():
        if isinstance(v, list):
            v = ", ".join(str(x) for x in v)
        ws_b.append([k, v])
    ws_b.column_dimensions["A"].width = 25
    ws_b.column_dimensions["B"].width = 60

    # Sauvegarde
    fname = _skill_filename(skill_meta, "xlsx")
    path = os.path.join(out_dir, fname)
    wb.save(path)
    return path


def export_skill_to_word(markdown: str, skill_meta: Dict[str, Any], out_dir: str) -> str:
    """Génère un .docx narratif à partir du markdown du skill."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    blocks = _parse_markdown(markdown)
    doc = Document()

    # Style global
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Page de garde
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(skill_meta.get("label", "Research").upper())
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0E, 0x4A, 0x2C)

    # Sub-title : answers compact
    answers = skill_meta.get("answers", {}) or {}
    sub_parts = []
    for key in ("ticker", "tickers", "trimestre", "annee_fiscale", "secteur", "univers",
                "style", "geo", "side", "horizon"):
        v = answers.get(key)
        if v:
            if isinstance(v, list):
                v = ", ".join(str(x) for x in v)
            sub_parts.append(str(v))
    if sub_parts:
        sub = doc.add_paragraph()
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = sub.add_run(" — ".join(sub_parts[:5]))
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(date.today().strftime("%d %B %Y"))
    run.font.size = Pt(11)
    run.font.italic = True
    doc.add_page_break()

    # Cadrage
    doc.add_heading("Paramètres de la recherche", level=1)
    for k, v in answers.items():
        if v in (None, "", []):
            continue
        if isinstance(v, list):
            v = ", ".join(str(x) for x in v)
        p = doc.add_paragraph()
        p.add_run(f"{k.replace('_', ' ').title()} : ").bold = True
        p.add_run(str(v))
    doc.add_paragraph()

    # Corps : on parcourt les blocs
    for btype, level, content in blocks:
        if btype == "h1":
            doc.add_heading(_strip_markdown_inline(content), level=1)
        elif btype == "h2":
            doc.add_heading(_strip_markdown_inline(content), level=2)
        elif btype == "h3":
            doc.add_heading(_strip_markdown_inline(content), level=3)
        elif btype == "para":
            doc.add_paragraph(_strip_markdown_inline(content))
        elif btype == "bullet":
            for item in content:
                doc.add_paragraph(_strip_markdown_inline(item), style="List Bullet")
        elif btype == "table" and content:
            n_rows = len(content)
            n_cols = max(len(r) for r in content) if content else 1
            tbl = doc.add_table(rows=n_rows, cols=n_cols)
            tbl.style = "Light Grid Accent 1"
            for ri, row_cells in enumerate(content):
                for ci, cell_val in enumerate(row_cells):
                    if ci < n_cols:
                        cell = tbl.cell(ri, ci)
                        cell.text = _strip_markdown_inline(cell_val)
                        if ri == 0:
                            for p in cell.paragraphs:
                                for r in p.runs:
                                    r.font.bold = True
            doc.add_paragraph()

    # Sauvegarde
    fname = _skill_filename(skill_meta, "docx")
    path = os.path.join(out_dir, fname)
    doc.save(path)
    return path


def export_skill_to_pptx(markdown: str, skill_meta: Dict[str, Any], out_dir: str) -> str:
    """Génère un .pptx pitch deck à partir du markdown du skill."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    blocks = _parse_markdown(markdown)
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_DARK = RGBColor(0x0A, 0x14, 0x18)
    ACCENT  = RGBColor(0x22, 0xC5, 0x5E)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    GREY    = RGBColor(0xAA, 0xAA, 0xAA)

    def _set_bg(slide, color=BG_DARK):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text(slide, left, top, width, height, text, size=18, bold=False,
                  color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        # Premier paragraphe
        first = True
        for paragraph_text in (text or "").split("\n"):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = paragraph_text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return tf

    blank = prs.slide_layouts[6]
    answers = skill_meta.get("answers", {}) or {}

    # ── Cover ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text(s, 0.5, 2.5, 12.3, 1.2, skill_meta.get("label", "Research").upper(),
              size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    sub_parts = []
    for key in ("ticker", "tickers", "trimestre", "annee_fiscale", "secteur", "univers",
                "style", "geo", "side", "horizon"):
        v = answers.get(key)
        if v:
            if isinstance(v, list):
                v = ", ".join(str(x) for x in v)
            sub_parts.append(str(v))
    if sub_parts:
        _add_text(s, 0.5, 3.9, 12.3, 0.8, " — ".join(sub_parts[:5]),
                  size=22, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, 0.5, 6.5, 12.3, 0.5, date.today().strftime("%d %B %Y"),
              size=14, color=GREY, align=PP_ALIGN.CENTER)

    # ── Slides par section H1 ────────────────────────────────────────────────
    # On regroupe les blocs en sections (chaque H1 démarre une nouvelle slide)
    sections: List[Tuple[str, List[Tuple[str, int, Any]]]] = []
    cur_title = "Introduction"
    cur_blocks: List[Tuple[str, int, Any]] = []
    for btype, level, content in blocks:
        if btype == "h1":
            if cur_blocks:
                sections.append((cur_title, cur_blocks))
            cur_title = _strip_markdown_inline(content)
            cur_blocks = []
        else:
            cur_blocks.append((btype, level, content))
    if cur_blocks:
        sections.append((cur_title, cur_blocks))

    # Si aucune section H1, créer une slide unique
    if not sections:
        sections = [("Synthèse", blocks)]

    for title, sec_blocks in sections:
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _add_text(s, 0.5, 0.4, 12.3, 0.8, title, size=28, bold=True, color=ACCENT)

        # Construit le contenu de la slide en concaténant
        body_lines: List[str] = []
        for btype, _, content in sec_blocks:
            if btype == "h2":
                body_lines.append("")
                body_lines.append(_strip_markdown_inline(content).upper())
            elif btype == "h3":
                body_lines.append(_strip_markdown_inline(content))
            elif btype == "para":
                txt = _strip_markdown_inline(content)
                # Tronque si trop long pour une slide
                if len(txt) > 350:
                    txt = txt[:347] + "..."
                body_lines.append(txt)
            elif btype == "bullet":
                for item in content[:6]:
                    body_lines.append("• " + _strip_markdown_inline(item))
            elif btype == "table" and content:
                # Petit récap : nb lignes/cols
                body_lines.append(f"[Tableau {len(content)} lignes × {len(content[0]) if content else 0} colonnes — voir Excel]")

        body_text = "\n".join(body_lines[:25])  # Cap à 25 lignes par slide
        _add_text(s, 0.5, 1.4, 12.3, 5.5, body_text, size=13, color=WHITE)

    # ── Slide finale : paramètres ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text(s, 0.5, 0.4, 12.3, 0.8, "Paramètres de la recherche",
              size=28, bold=True, color=ACCENT)
    lines = []
    for k, v in answers.items():
        if v in (None, "", []):
            continue
        if isinstance(v, list):
            v = ", ".join(str(x) for x in v)
        lines.append(f"{k.replace('_', ' ').title()} : {v}")
    _add_text(s, 0.5, 1.4, 12.3, 5.5, "\n".join(lines), size=14, color=WHITE)

    # Sauvegarde
    fname = _skill_filename(skill_meta, "pptx")
    path = os.path.join(out_dir, fname)
    prs.save(path)
    return path
